import sys
import tempfile
from pathlib import Path
from collections import defaultdict
import numpy as np
import pandas as pd
import streamlit as st
from sklearn.metrics.pairwise import cosine_similarity
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.docstore.document import Document
import altair as alt

# ---------------------- 백엔드 함수 정의 ----------------------

def load_documents_from_file(file_path: str) -> list[Document]:
    """pdf, docx, txt, xlsx, ppt, pptx 등 다양한 포맷에서 텍스트를 로드하여 Document 리스트로 반환합니다."""
    suffix = Path(file_path).suffix.lower()
    docs: list[Document] = []

    if suffix == ".pdf":
        docs = PyPDFLoader(file_path).load()
    elif suffix in {".docx", ".doc"}:
        docs = Docx2txtLoader(file_path).load()
    elif suffix == ".txt":
        docs = TextLoader(file_path).load()
    elif suffix in {".xlsx", ".xls"}:
        xls = pd.read_excel(file_path, sheet_name=None)
        text = []
        for sheet_name, df in xls.items():
            text.append(f"### Sheet: {sheet_name}")
            text.append(df.astype(str).to_csv(index=False, sep=" "))
        docs = [Document(page_content="\n".join(text), metadata={"source": file_path})]
    elif suffix in {".pptx", ".ppt"}:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            docs = [Document(page_content="\n".join(text), metadata={"source": file_path})]
        except ImportError:
            docs = []
    else:
        raise ValueError(f"지원되지 않는 파일 포맷: {suffix}")

    return docs


def compute_similarity(uploaded_path: str, db_path: str) -> list[tuple[str, float]]:
    docs = load_documents_from_file(uploaded_path)
    if not docs:
        raise ValueError(f"업로드된 파일에 내용이 없습니다: {uploaded_path}")

    query_text = docs[0].page_content
    embedder = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    query_vec = np.array(embedder.embed_documents([query_text])[0]).reshape(1, -1)

    vectordb = FAISS.load_local(db_path, embedder, allow_dangerous_deserialization=True)
    index = vectordb.index
    total = index.ntotal
    stored_vecs = index.reconstruct_n(0, total)
    stored_matrix = np.array(stored_vecs)
    meta_docs = list(vectordb.docstore._dict.values())

    sims = cosine_similarity(query_vec, stored_matrix)[0]
    file_scores = defaultdict(float)
    for md, sim in zip(meta_docs, sims):
        src = md.metadata.get("source", "")
        fname = Path(src).name if src else "알 수 없음"
        file_scores[fname] = max(file_scores[fname], sim)

    results = sorted(file_scores.items(), key=lambda x: -x[1])
    return [(fn, f"{int(round(sc * 100))}%") for fn, sc in results]


def rebuild_db(source_folder: str, db_path: str):
    embedder = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    docs: list[Document] = []
    for file in Path(source_folder).iterdir():
        try:
            docs.extend(load_documents_from_file(str(file)))
        except Exception:
            continue
    if not docs:
        raise ValueError(f"지원되는 파일이 없습니다: {source_folder}")
    vectordb = FAISS.from_documents(docs, embedder)
    vectordb.save_local(db_path)

# ---------------------- Streamlit UI 설정 ----------------------
st.set_page_config(
    page_title="document Analyzer", 
    layout="wide", 
    initial_sidebar_state="expanded")

# CSS 커스터마이징
st.markdown(
    """
    <style>
    /* 버튼을 입력칸 크기와 맞추고, 좌측 정렬 */
    div.stButton > button {
        width: 100% !important;              /* 전체 폭 */
        background-color: #1f77b4 !important;
        color: white !important;
        font-weight: bold;
        border: none;
        border-radius: 6px;
        padding: 0.5rem 1rem;
        text-align: left !important;         /* 텍스트 좌측 정렬 */
        transition: background-color 0.3s ease;
    }

    div.stButton > button:hover {
        background-color: #135b91 !important;
    }

    /* 전체 앱 기본 폰트 스타일 */
    .stApp {
        font-size: 18px;
        font-family: 'Apple SD Gothic Neo', 'Malgun Gothic', sans-serif;
    }

    /* 좌측 사이드바 글자 크기 */
    section[data-testid="stSidebar"] * {
        font-size: 1.3rem !important;
        font-weight: 500 !important;
    }

    /* 사이드바 전체 컴포넌트 간 여백 최소화 */
    section[data-testid="stSidebar"] div {
        margin-top: 0.1rem !important;
        margin-bottom: 0.1rem !important;
    }

    /* 사이드바 영역 폭 제한 */
    .css-1x92ibr {
        max-width: 90%;
        margin: auto;
    }

    /* 사이드바 입력 필드 여백 축소 */
    div[data-testid="stSidebar"] input {
        padding-top: 0.3rem !important;
        padding-bottom: 0.3rem !important;
    }

    /* 버튼 크기 및 여백 조정 */
    .stButton > button {
        font-size: 1.2rem !important;
        padding: 0.5rem 1rem !important;
    }

    /* 파일 업로더 레이블 스타일 */
    .stFileUploader label {
        font-size: 1.3rem !important;
        font-weight: 600 !important;
    }

    /* 좌우 헤더(h2) 통일 */
    h2, .st-emotion-cache-1kyxreq {
        font-size: 2rem !important;
        font-weight: 700 !important;
    }
    </style>
    """, unsafe_allow_html=True
)


st.title("📄 document Analyzer")
sidebar, main = st.columns([1, 4])

with sidebar:
    # HTML 헤더
    st.markdown("<div style='font-size:2rem; font-weight:bold; margin-top:1rem;'>Index Folder</div>", 
                unsafe_allow_html=True)

    # HTML 레이블만 찍고, 입력 박스는 label 없이
    st.markdown("<div style='font-size:1rem; font-weight:bold; margin-bottom:0.25rem;'>Source Folder</div>",
                unsafe_allow_html=True)
    src_folder = st.text_input("Source Folder", value="source_documents", label_visibility="collapsed")

    st.markdown("<div style='font-size:1rem; font-weight:bold; margin-bottom:0.25rem;'>FAISS Index Folder</div>",
                unsafe_allow_html=True)
    db_folder = st.text_input("FAISS Index Folder", value="db", label_visibility="collapsed")

    # 레이블만 HTML로 찍어두고, 실제 버튼은 빈 텍스트
    st.markdown("<div style='font-size:1rem; font-weight:bold; margin-top:1rem;'>Index Rebuilder</div>",
                unsafe_allow_html=True)
   # db_folder = st.text_input("", value="db")
    
    if st.button("Start"):
        try:
            rebuild_db(src_folder, db_folder)
            st.success(f"인덱스 재빌드 완료: {src_folder} → {db_folder}")
        except Exception as e:
            st.error(f"인덱스 재빌드 실패: {e}")

with main:
    st.header("File Uploader")
    uploaded_file = st.file_uploader(
        "파일을 드래그 앤 드롭하세요",
        type=["pdf", "doc", "docx", "txt", "xlsx", "ppt", "pptx"],
        label_visibility="visible"
    )

    if uploaded_file:
        suffix = Path(uploaded_file.name).suffix.lower()
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        tmp.write(uploaded_file.getbuffer()); tmp.flush()

        st.info(f"업로드됨: **{uploaded_file.name}**")
        st.info(f"사용 중인 DB: **{db_folder}**")

        try:
            results = compute_similarity(tmp.name, db_folder)
            df = pd.DataFrame(results, columns=["파일명", "유사도 (%)"])
            df["유사도"] = df["유사도 (%)"].str.replace('%', '').astype(int)
           # .sort_values(by="유사도 (%)", ascending=False)

            # ✅ 여기에 이 코드 삽입
            df_display = df[["파일명", "유사도 (%)"]]
            styled_df = df_display.style.set_table_styles([
                {'selector': 'th', 'props': [('background-color', '#4B8BBE'), ('color', 'white')]}
            ])

            st.subheader("모든 파일의 유사도 (%)")
            st.dataframe(df[["파일명", "유사도 (%)"]], use_container_width=True)

            st.subheader("유사도 가로 막대 차트")
            chart = alt.Chart(df).mark_bar(cornerRadius=5).encode(
                y=alt.Y('파일명:N', sort='-x', title=None, axis=alt.Axis(labelFontSize=14)),
                x=alt.X('유사도:Q', title='유사도 (%)'),
                color=alt.value('#4B8BBE')
            ).properties(height=400)
            st.altair_chart(chart, use_container_width=True)

            df["row"] = " "  # 단일 행 구성

            heatmap = alt.Chart(df).mark_rect().encode(
                x=alt.X('파일명:N', axis=alt.Axis(labelAngle=-45, labelFontSize=12)),
                y=alt.Y('row:O', axis=None),
                color=alt.Color('유사도:Q', scale=alt.Scale(scheme='viridis'), title='유사도 (%)')
            ).properties(height=80)

            st.altair_chart(heatmap, use_container_width=True)

        except Exception as e:
            st.error(f"오류 발생: {e}")


