import sys
import tempfile
from pathlib import Path
from collections import defaultdict
import numpy as np
import pandas as pd
import streamlit as st
from sklearn.metrics.pairwise import cosine_similarity
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.docstore.document import Document
import altair as alt

# ---------------------- 백엔드 함수 정의 ----------------------

def load_documents_from_file(file_path: str) -> list[Document]:
    """pdf, docx, txt, xlsx, ppt, pptx 등 다양한 포맷에서 텍스트를 로드하여 Document 리스트로 반환합니다."""
    suffix = Path(file_path).suffix.lower()
    docs: list[Document] = []

    if suffix == ".pdf":
        docs = PyPDFLoader(file_path).load()
    elif suffix in {".docx", ".doc"}:
        docs = Docx2txtLoader(file_path).load()
    elif suffix == ".txt":
        docs = TextLoader(file_path).load()
    elif suffix in {".xlsx", ".xls"}:
        xls = pd.read_excel(file_path, sheet_name=None)
        text = []
        for sheet_name, df in xls.items():
            text.append(f"### Sheet: {sheet_name}")
            text.append(df.astype(str).to_csv(index=False, sep=" "))
        docs = [Document(page_content="\n".join(text), metadata={"source": file_path})]
    elif suffix in {".pptx", ".ppt"}:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            docs = [Document(page_content="\n".join(text), metadata={"source": file_path})]
        except ImportError:
            docs = []
    else:
        raise ValueError(f"지원되지 않는 파일 포맷: {suffix}")

    return docs


def compute_similarity(uploaded_path: str, db_path: str) -> list[tuple[str, float]]:
    docs = load_documents_from_file(uploaded_path)
    if not docs:
        raise ValueError(f"업로드된 파일에 내용이 없습니다: {uploaded_path}")

    query_text = docs[0].page_content
    embedder = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    query_vec = np.array(embedder.embed_documents([query_text])[0]).reshape(1, -1)

    vectordb = FAISS.load_local(db_path, embedder, allow_dangerous_deserialization=True)
    index = vectordb.index
    total = index.ntotal
    stored_vecs = index.reconstruct_n(0, total)
    stored_matrix = np.array(stored_vecs)
    meta_docs = list(vectordb.docstore._dict.values())

    sims = cosine_similarity(query_vec, stored_matrix)[0]
    file_scores = defaultdict(float)
    for md, sim in zip(meta_docs, sims):
        src = md.metadata.get("source", "")
        fname = Path(src).name if src else "알 수 없음"
        file_scores[fname] = max(file_scores[fname], sim)

    results = sorted(file_scores.items(), key=lambda x: -x[1])
    return [(fn, sc * 100) for fn, sc in results]


def rebuild_db(source_folder: str, db_path: str):
    embedder = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    docs: list[Document] = []
    for file in Path(source_folder).iterdir():
        try:
            docs.extend(load_documents_from_file(str(file)))
        except Exception:
            continue
    if not docs:
        raise ValueError(f"지원되는 파일이 없습니다: {source_folder}")
    vectordb = FAISS.from_documents(docs, embedder)
    vectordb.save_local(db_path)

# ---------------------- Streamlit UI 설정 ----------------------
st.set_page_config(page_title="document Analyzer", layout="wide", initial_sidebar_state="expanded")

# CSS 커스터마이징
st.markdown(
    """
    <style>
    /* 사이드바 전체 글자 200% */
    section[data-testid="stSidebar"] * { font-size: 200% !important; }
    .css-1x92ibr { max-width: 90%; margin: auto; }
    
    /* 사이드바 내부의 모든 레이블 텍스트를 200%로 키우고 볼드 처리 */
    section[data-testid="stSidebar"] label,
    section[data-testid="stSidebar"] .stButton>button {
        font-size: 400% !important;
        font-weight: bold !important;
    }

    /* 메인 영역의 header (h2) 크기도 더 크게 조정 */
    .css-1lsmgbg h2, /* 클래스명은 버전에 따라 다를 수 있으니, 개발자 도구로 확인하세요 */
    .css-1lsmgbg .stHeader {
        font-size: 4.5rem !important;
    }
    
    </style>
    """, unsafe_allow_html=True
)

st.title("📄 document Analyzer")
sidebar, main = st.columns([1, 4])

with sidebar:
    st.header("Index Folder")
    src_folder = st.text_input("소스 폴더", value="source_documents")
    db_folder = st.text_input("FAISS 인덱스 폴더", value="db")
    if st.button("인덱스 재빌드"):
        try:
            rebuild_db(src_folder, db_folder)
            st.success(f"인덱스 재빌드 완료: {src_folder} → {db_folder}")
        except Exception as e:
            st.error(f"인덱스 재빌드 실패: {e}")

with main:
    st.header("File Uploader")
    uploaded_file = st.file_uploader(
        "파일을 드래그 앤 드롭하세요",
        type=["pdf", "doc", "docx", "txt", "xlsx", "ppt", "pptx"],
        label_visibility="visible"
    )

    if uploaded_file:
        suffix = Path(uploaded_file.name).suffix.lower()
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        tmp.write(uploaded_file.getbuffer()); tmp.flush()

        st.info(f"업로드됨: **{uploaded_file.name}**")
        st.info(f"사용 중인 DB: **{db_folder}**")

        try:
            results = compute_similarity(tmp.name, db_folder)
            df = pd.DataFrame(results, columns=["파일명", "유사도 (%)"]).sort_values(by="유사도 (%)", ascending=False)

            st.subheader("모든 파일의 유사도 (%)")
            st.dataframe(df, use_container_width=True)

            st.subheader("유사도 가로 막대 차트")
            chart = alt.Chart(df).mark_bar(cornerRadius=5).encode(
                y=alt.Y('파일명:N', sort='-x', title=None, axis=alt.Axis(labelFontSize=14)),
                x=alt.X('유사도 (%):Q', title='유사도 (%)'),
                color=alt.value('#4B8BBE')
            ).properties(height=400)
            st.altair_chart(chart, use_container_width=True)

            df['row'] = ' '

            st.subheader("유사도 히트맵")
            heatmap = alt.Chart(df).mark_rect().encode(
                x=alt.X('파일명:N', sort='-row', axis=alt.Axis(labelAngle=-45, labelFontSize=12)),
                y=alt.Y('row:O', axis=None),
                color=alt.Color('유사도 (%):Q', scale=alt.Scale(scheme='viridis'), title='유사도 (%)')
            ).properties(height=80)
            st.altair_chart(heatmap, use_container_width=True)

        except Exception as e:
            st.error(f"오류 발생: {e}")


